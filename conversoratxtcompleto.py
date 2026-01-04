# conversoratxtcompleto.py
#Este es un conversor de archivos PDF, DOC, DOCX, PPT, PPTX a TXT
#A comparación de los otros programas de la suite, este no tiene una interfaz de inicio,
#simplemente pide que selecciones archivos (o toma los previamente seleccionados desde el menú contextual cuando este programa es ejecutado desde el menú contextual)
#y transforma eficientemente a txt.
#Solo aplica OCR de 75 dpi a archivos que tengan más de 10% de páginas con menos de 20 caracteres cada una.

from __future__ import annotations

import sys
import os
from concurrent.futures import ThreadPoolExecutor, wait, FIRST_COMPLETED
import threading
import queue
import shutil
import subprocess
import tempfile
from pathlib import Path
import tkinter as tk
from tkinter import filedialog
from typing import List

# Base
from pypdf import PdfReader
try:
    import fitz  # type: ignore
    _HAS_FITZ = True
except Exception:
    fitz = None  # type: ignore
    _HAS_FITZ = False
try:
    from ebooklib import epub  # type: ignore
    from ebooklib.epub import EpubHtml  # type: ignore
    _HAS_EPUB = True
except Exception:
    epub = None  # type: ignore
    EpubHtml = object  # type: ignore
    _HAS_EPUB = False

try:
    from bs4 import BeautifulSoup  # type: ignore
    _HAS_BS4 = True
except Exception:
    BeautifulSoup = None  # type: ignore
    _HAS_BS4 = False
from config import get_config, window_geometry, apply_theme, ProgressDialog
CONFIG = get_config('conversor')

# UI de progreso
def _norm_txt(s: str) -> str:
    return (s or "").replace("\r\n", "\n").replace("\r", "\n")


def _which(names: List[str]) -> str | None:
    for n in names:
        p = shutil.which(n)
        if p:
            return p
    return None



def _find_soffice() -> str | None:
    """Encuentra LibreOffice (soffice) incluso si no está en PATH."""
    p = _which(["soffice", "soffice.com", "lowriter", "soffice.bin"])
    if p:
        return p

    # Rutas típicas en Windows
    try:
        candidates: list[Path] = []
        for env in ("PROGRAMFILES", "PROGRAMFILES(X86)", "LOCALAPPDATA"):
            base = os.environ.get(env)
            if base:
                candidates.append(Path(base))
        for base in candidates:
            for rel in (
                r"LibreOffice\program\soffice.com",
                r"LibreOffice\program\soffice.exe",
                r"LibreOffice\program\soffice.bin",
            ):
                c = base / rel
                try:
                    if c.exists():
                        return str(c)
                except Exception:
                    pass
    except Exception:
        pass
    return None


def _pages_to_ranges(pages_1based: List[int]) -> str:
    """Convierte [1,2,3,7,8] -> '1-3,7-8' (formato OCRmyPDF --pages)."""
    try:
        pages = sorted({int(p) for p in pages_1based if int(p) >= 1})
    except Exception:
        return ""
    if not pages:
        return ""
    out: List[str] = []
    start = prev = pages[0]
    for p in pages[1:]:
        if p == prev + 1:
            prev = p
            continue
        out.append(f"{start}-{prev}" if start != prev else f"{start}")
        start = prev = p
    out.append(f"{start}-{prev}" if start != prev else f"{start}")
    return ",".join(out)

import contextlib
import io

@contextlib.contextmanager
def _silence_stdio():
    """Silencia stdout/stderr temporalmente (evita spam de dependencias como qpdf/pikepdf)."""
    try:
        with contextlib.redirect_stdout(io.StringIO()), contextlib.redirect_stderr(io.StringIO()):
            yield
    except Exception:
        # Si falla el redirect (muy raro), no bloquear la conversión.
        yield



class _PDFTextExtractor:
    """Hilo reutilizable para extraer texto de páginas PDF con timeout sin crear un hilo por página.
    Nota: si una extracción se cuelga, se crea una instancia nueva y se deja la anterior como daemon."""
    def __init__(self):
        self._q: "queue.Queue[tuple]" = queue.Queue()
        self._t = threading.Thread(target=self._loop, daemon=True)
        self._t.start()

    def _loop(self):
        while True:
            page_obj, out, ev = self._q.get()
            try:
                out["text"] = page_obj.extract_text() or ""
            except Exception:
                out["text"] = ""
            finally:
                try:
                    ev.set()
                except Exception:
                    pass

    def extract(self, page_obj, timeout_s: float) -> tuple[bool, str]:
        out: dict = {}
        ev = threading.Event()
        try:
            self._q.put((page_obj, out, ev))
        except Exception:
            return True, ""
        if ev.wait(timeout_s):
            return True, str(out.get("text", "") or "")
        return False, ""






class _PDFTextExtractorPool:
    """Pool de hilos *daemon* reutilizables para extracción PDF con timeout.
    Permite paralelizar extracción sin crear un hilo por página."""

    def __init__(self, workers: int):
        self.workers = max(1, int(workers) if workers else 1)
        self._lock = threading.Lock()
        self._rr = 0
        self._extractors = [_PDFTextExtractor() for _ in range(self.workers)]

    def extract(self, page_obj, timeout_s: float) -> tuple[bool, str]:
        # Selección round-robin (barata y suficiente).
        with self._lock:
            idx = self._rr
            self._rr = (self._rr + 1) % self.workers
            ex = self._extractors[idx]

        ok, txt = ex.extract(page_obj, timeout_s)
        if ok:
            return True, txt

        # Timeout: reemplazar extractor para evitar que el hilo colgado bloquee futuros.
        try:
            with self._lock:
                self._extractors[idx] = _PDFTextExtractor()
        except Exception:
            pass
        return False, ""


def _fast_pdf_pagecount(path: Path) -> int:
    """Conteo de páginas para progreso global.

    Prioridad:
      1) PyMuPDF (rápido y robusto).
      2) Escaneo de /Type /Page (fallback muy barato, no 100% exacto).
    """
    # 1) PyMuPDF
    if _HAS_FITZ and (fitz is not None):
        try:
            doc = fitz.open(str(path))
            try:
                pc = int(getattr(doc, 'page_count', 0) or 0)
                if pc > 0:
                    return pc
            finally:
                try:
                    doc.close()
                except Exception:
                    pass
        except Exception:
            pass

    # 2) Fallback: escaneo por ocurrencias de /Type /Page
    try:
        import re
        pat = re.compile(br"/Type\s*/Page\b")
        count = 0
        tail = b""
        with path.open('rb') as f:
            while True:
                chunk = f.read(1024 * 1024)
                if not chunk:
                    break
                data = tail + chunk
                count += len(pat.findall(data))
                tail = data[-64:]
        return int(count)
    except Exception:
        return 0


def _estimate_units_for_path(path: Path) -> int:
    """Estimación rápida para un progreso GLOBAL estable (toda la lista)."""
    ext = path.suffix.lower()
    try:
        if ext == ".pdf":
            pc = _fast_pdf_pagecount(path)
            return max(1, pc)
        if ext == ".epub":
            if (not _HAS_EPUB) or (epub is None):
                return 1
            book = epub.read_epub(str(path))
            items = [it for it in book.get_items() if isinstance(it, EpubHtml)]
            return max(1, len(items))
        if ext == ".docx":
            import zipfile
            with zipfile.ZipFile(path) as z:
                xml = z.read("word/document.xml")
            return max(1, xml.count(b"<w:p"))
        if ext == ".pptx":
            import zipfile
            with zipfile.ZipFile(path) as z:
                xml = z.read("ppt/presentation.xml")
            return max(1, xml.count(b"<p:sldId"))
    except Exception:
        pass
    return 1


class HeadlessConverterApp:
    def __init__(self, files: List[str]):
        self._root = tk.Tk()
        # ProgressDialog ahora es overlay (misma ventana). No se oculta el root.
        # ProgressDialog ahora es overlay (misma ventana).
        # Se configura el root con el mismo "window_geometry" que el resto de apps (evita ventana mini).
        try:
            self._root.withdraw()
        except Exception:
            pass

        pconf = get_config("progress")
        try:
            sw, sh = self._root.winfo_screenwidth(), self._root.winfo_screenheight()

            # IMPORTANTE:
            # El conversor completo no tiene UI "grande" (solo progreso), pero el usuario
            # espera que la ventana tenga las mismas dimensiones que el resto de apps.
            # Por eso el root usa el geometry del app "conversor" (no el del overlay "progress").
            ww, wh = window_geometry("conversor", sw, sh)

            # Centrado real
            x = max(0, int((sw - ww) / 2))
            y = max(0, int((sh - wh) / 2))
            self._root.geometry(f"{ww}x{wh}+{x}+{y}")

            # Min size (coherente con el resto de apps)
            try:
                self._root.minsize(int(CONFIG["window"]["min_width"]), int(CONFIG["window"]["min_height"]))
            except Exception:
                self._root.minsize(int(pconf["window"]["min_width"]), int(pconf["window"]["min_height"]))

            # Tema / icono coherente con el resto de la suite
            try:
                # Nota: aplicamos tokens del app "conversor" (tamaño/estilo de ventana),
                # pero el título/ícono siguen siendo los del modo progreso.
                apply_theme(self._root, app="conversor", title=pconf["window"]["title"], icon=pconf["window"].get("icon_name"))
            except Exception:
                # Fallback mínimo
                try:
                    self._root.title(pconf["window"]["title"])
                except Exception:
                    pass

            try:
                self._root.configure(bg=pconf["colors"]["app_bg"])
            except Exception:
                pass
        except Exception:
            pass

        try:
            # Permite cerrar con X sin dejar procesos colgados
            self._root.protocol("WM_DELETE_WINDOW", self._final_close)
        except Exception:
            pass

        self._root.deiconify()

        self._icon_dir = Path(__file__).parent
        self._queue: "queue.Queue[tuple]" = queue.Queue()

        # Progreso GLOBAL (toda la lista), no por-archivo
        self._g_total_units: int = 0
        self._g_done_units: int = 0
        self._g_pages_total: int = 0  # conteo lógico (sin OCR extra)
        self._g_pages_done: int = 0
        self._cur_file_name: str = ""
        self._last_file_unit: int = 0
        self._last_ocr_unit: int = 0
        self._files = [Path(f) for f in files] if files else []
        if not self._files:
            paths = filedialog.askopenfilenames(
                title="Selecciona archivos",
                filetypes=[
                    ("Todos soportados", "*.pdf;*.epub;*.doc;*.docx;*.ppt;*.pptx"),
                    ("PDF", "*.pdf"),
                    ("EPUB", "*.epub"),
                    ("Word", "*.doc;*.docx"),
                    ("PowerPoint", "*.ppt;*.pptx"),
                    ("Todos", "*.*"),
                ]
            )
            self._files = [Path(p) for p in paths]

        if not self._files:
            self._root.destroy()
            return

        try:
            self._pdlg = ProgressDialog(parent=self._root, icon_search_dir=self._icon_dir)
        except TypeError:
            self._pdlg = ProgressDialog(parent=self._root)
        except Exception:
            self._pdlg = ProgressDialog(self._root)
        try:
            self._pdlg.set_title("Preparando…")
            self._pdlg.set_subtitle("Inicializando…")
            self._pdlg.set_progress(0.0)
            self._pdlg.show()
            pass  # update_idletasks deshabilitado por estabilidad (Tk 9 / Python 3.14)
        except Exception:
            pass
        _target = getattr(self, '_process_files', None)
        if not callable(_target):
            _target = getattr(self, '_process_file', None)
        if not callable(_target):
            def _target():
                self._emit('error', 'Error interno: falta _process_files en este build.')
                self._emit('done')

        self._worker = threading.Thread(target=_target, daemon=True)
        self._worker.start()
        self._root.after(60, self._pump_progress)
        self._root.mainloop()


    def _emit(self, tag: str, *data) -> None:
        """Encola eventos desde el hilo worker hacia el UI thread."""
        try:
            self._queue.put((tag, *data))
        except Exception:
            pass

    # ---------------- PDF ----------------
    def _safe_extract_pdf_page(self, page):
        """Extrae texto de una página PDF con protección anti-bloqueo.
        Optimización: reusa un único hilo extractor; solo crea otro si alguna página se cuelga."""
        try:
            timeout = float(CONFIG.get("app", {}).get("safe_extract_timeout", 6.0))
        except Exception:
            timeout = 6.0

        # Si el usuario desactiva el timeout (<=0), no hay sobrecosto de hilos.
        if timeout <= 0:
            try:
                return page.extract_text() or ""
            except Exception:
                return ""

        try:
            # Pool de extractores para permitir paralelismo con timeout real.
            try:
                workers = int((CONFIG.get("app", {}) or {}).get("pdf_text_workers", 0) or 0)
            except Exception:
                workers = 0
            if not workers:
                workers = max(1, min(4, (os.cpu_count() or 4)))

            if (not hasattr(self, "_pdf_extractor_pool")) or (self._pdf_extractor_pool is None) or (getattr(self._pdf_extractor_pool, "workers", None) != workers):
                self._pdf_extractor_pool = _PDFTextExtractorPool(workers)

            ok, txt = self._pdf_extractor_pool.extract(page, timeout)
            if ok:
                return txt
            return ""
        except Exception:

            return ""

    # ---------------- DOCX (python-docx -> fallback ZIP) ----------------
    def _convert_docx(self, path: Path):
        # 1) intentar python-docx si está disponible
        try:
            from docx import Document  # type: ignore
            try:
                doc = Document(str(path))
            except Exception as e:
                self._emit("error", f"No se pudo abrir .docx (python-docx): {e}")
                text = self._docx_from_zip(path)
                return self._write_simple(text, path)
            chunks: List[str] = []
            total = max(1, len(doc.paragraphs) + sum(len(t.rows) for t in getattr(doc, "tables", [])))
            done = 0
            for p in doc.paragraphs:
                done += 1; self._emit("extract_page", done, total)
                chunks.append(p.text or "")
            for t in getattr(doc, "tables", []):
                for row in t.rows:
                    cells = [c.text or "" for c in row.cells]
                    if cells:
                        chunks.append("\t".join(cells))
                    done += 1; self._emit("extract_page", min(done, total), total)
            text = _norm_txt("\n".join(chunks))
            return self._write_simple(text, path)
        except Exception:
            # 2) fallback puro Python
            text = self._docx_from_zip(path)
            return self._write_simple(text, path)

    def _docx_from_zip(self, path: Path) -> str:
        """Extrae texto de DOCX via zip + XML sin dependencias."""
        import zipfile, xml.etree.ElementTree as ET
        ns = {"w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main"}
        def extract_file(zf, name):
            try:
                xml = zf.read(name)
                root = ET.fromstring(xml)
                parts: List[str] = []
                for p in root.iterfind(".//w:p", ns):
                    runs = []
                    for t in p.iterfind(".//w:t", ns):
                        if t.text:
                            runs.append(t.text)
                    parts.append("".join(runs))
                return "\n".join(parts)
            except Exception:
                return ""
        texts: List[str] = []
        with zipfile.ZipFile(str(path), "r") as zf:
            texts.append(extract_file(zf, "word/document.xml"))
            for i in range(1, 10):
                for kind in ("header", "footer"):
                    name = f"word/{kind}{i}.xml"
                    if name in zf.namelist():
                        texts.append(extract_file(zf, name))
        return _norm_txt("\n".join([t for t in texts if t]))

    # ---------------- PPTX (python-pptx -> fallback ZIP) ----------------
    def _convert_pptx(self, path: Path):
        # 1) intentar python-pptx
        try:
            from pptx import Presentation  # type: ignore
            try:
                prs = Presentation(str(path))
            except Exception as e:
                self._emit("error", f"No se pudo abrir .pptx (python-pptx): {e}")
                text = self._pptx_from_zip(path)
                return self._write_simple(text, path)
            chunks: List[str] = []
            total = max(1, len(prs.slides))
            def shape_texts(shape):
                out = []
                try:
                    if getattr(shape, "has_text_frame", False) and shape.text_frame:
                        out.append(shape.text_frame.text or "")
                    if hasattr(shape, "table") and shape.table is not None:
                        tbl = shape.table
                        for r in tbl.rows:
                            cells = []
                            for c in r.cells:
                                try:
                                    cells.append(c.text or "")
                                except Exception:
                                    cells.append("")
                            if cells:
                                out.append("\t".join(cells))
                    if hasattr(shape, "shapes"):
                        for sh in shape.shapes:
                            out.extend(shape_texts(sh))
                except Exception:
                    pass
                return out
            for i, slide in enumerate(prs.slides, 1):
                texts = []
                for sh in slide.shapes:
                    texts.extend(shape_texts(sh))
                chunks.append(f"[Diapositiva {i}]\n" + "\n".join(texts).strip())
                self._emit("extract_page", i, total)
            return self._write_simple(_norm_txt("\n\n".join(chunks)), path)
        except Exception:
            # 2) fallback puro Python
            text = self._pptx_from_zip(path)
            return self._write_simple(text, path)

    def _pptx_from_zip(self, path: Path) -> str:
        """Extrae texto de PPTX via zip + XML sin dependencias."""
        import zipfile, re, xml.etree.ElementTree as ET
        ns = {
            "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
            "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
        }
        texts: List[str] = []
        with zipfile.ZipFile(str(path), "r") as zf:
            slides = [n for n in zf.namelist() if n.startswith("ppt/slides/slide") and n.endswith(".xml")]
            def slide_index(name: str) -> int:
                import re
                m = re.search(r"slide(\d+)\.xml", name)
                return int(m.group(1)) if m else 0
            slides.sort(key=slide_index)
            for idx, name in enumerate(slides, 1):
                try:
                    xml = zf.read(name)
                    root = ET.fromstring(xml)
                    ts = []
                    for t in root.iter("{%s}t" % ns["a"]):
                        if t.text:
                            ts.append(t.text)
                    slide_txt = "\n".join(ts).strip()
                    texts.append(f"[Diapositiva {idx}]\n{slide_txt}")
                except Exception:
                    texts.append(f"[Diapositiva {idx}]\n")
        return _norm_txt("\n\n".join(texts))

    # ---------------- DOC (COM -> LibreOffice) ----------------
    def _convert_doc(self, path: Path):
        # 1) Word por COM
        try:
            import pythoncom  # type: ignore
            pythoncom.CoInitialize()
            import win32com.client  # type: ignore
            try:
                wdFormatUnicodeText = 7
                try:
                    word = win32com.client.gencache.EnsureDispatch("Word.Application")
                except Exception:
                    word = win32com.client.Dispatch("Word.Application")
                word.Visible = False
                word.DisplayAlerts = 0
                doc = word.Documents.Open(str(path))
                out_file = path.with_suffix(".txt")
                doc.SaveAs(str(out_file), FileFormat=wdFormatUnicodeText)
                doc.Close(False)
                word.Quit()
                self._emit("saved", out_file.name)
                try: pythoncom.CoUninitialize()
                except Exception: pass
                return
            except Exception as e:
                # si Word falla, liberamos COM y pasamos al plan B
                try: pythoncom.CoUninitialize()
                except Exception: pass
                self._emit("error", f"Word/COM no disponible, probando LibreOffice… ({e})")
        except Exception:
            # no hay COM/pywin32, ir al plan B
            self._emit("error", "Word/COM no disponible, probando LibreOffice…")

        # 2) Fallback LibreOffice headless
        if self._fallback_doc_via_libreoffice(path):
            return

        # 3) Sin salida
        self._emit("error", "Para .doc necesitas Word+pywin32 o LibreOffice instalado para conversión headless.")

    def _fallback_doc_via_libreoffice(self, path: Path) -> bool:
        """
        Intenta convertir .doc usando LibreOffice:
        - Primero a DOCX, luego reusa el extractor DOCX.
        - Si falla, intenta exportar a TXT directo.
        """
        soffice = _find_soffice()
        if not soffice:
            return False

        # A) Convertir a DOCX
        try:
            with tempfile.TemporaryDirectory() as td:
                outdir = Path(td)
                proc = subprocess.run(
                    [soffice, "--headless", "--convert-to", "docx", "--outdir", str(outdir), str(path)],
                    capture_output=True, text=True
                )
                if proc.returncode == 0:
                    conv = outdir / (path.stem + ".docx")
                    if conv.exists():
                        # usar el lector DOCX puro
                        text = self._docx_from_zip(conv)
                        self._write_simple(text, path)
                        return True
        except Exception:
            pass

        # B) Convertir a TXT directo
        try:
            with tempfile.TemporaryDirectory() as td:
                outdir = Path(td)
                proc = subprocess.run(
                    [soffice, "--headless", "--convert-to", "txt:Text", "--outdir", str(outdir), str(path)],
                    capture_output=True, text=True
                )
                if proc.returncode == 0:
                    conv = outdir / (path.stem + ".txt")
                    if conv.exists():
                        txt = conv.read_text(encoding="utf-8", errors="ignore")
                        self._write_simple(_norm_txt(txt), path)
                        return True
        except Exception:
            pass
        return False

    
    def _fallback_ppt_via_libreoffice(self, path: Path) -> bool:
        """Intenta convertir .ppt usando LibreOffice headless.

        Estrategia:
        - Convertir a PPTX y extraer texto con el parser PPTX (ZIP+XML).
        - Si falla, intentar exportar TXT directo.
        """
        soffice = _find_soffice()
        if not soffice:
            return False

        # A) Convertir a PPTX (preferido)
        try:
            with tempfile.TemporaryDirectory() as td:
                outdir = Path(td)
                proc = subprocess.run(
                    [soffice, "--headless", "--convert-to", "pptx", "--outdir", str(outdir), str(path)],
                    capture_output=True, text=True
                )
                if proc.returncode == 0:
                    conv = outdir / (path.stem + ".pptx")
                    if conv.exists():
                        text = self._pptx_from_zip(conv)
                        self._write_simple(text, path)
                        return True
        except Exception:
            pass

        # B) Convertir a TXT directo
        try:
            with tempfile.TemporaryDirectory() as td:
                outdir = Path(td)
                proc = subprocess.run(
                    [soffice, "--headless", "--convert-to", "txt:Text", "--outdir", str(outdir), str(path)],
                    capture_output=True, text=True
                )
                if proc.returncode == 0:
                    conv = outdir / (path.stem + ".txt")
                    if conv.exists():
                        txt = conv.read_text(encoding="utf-8", errors="ignore")
                        self._write_simple(_norm_txt(txt), path)
                        return True
        except Exception:
            pass

        return False


# ---------------- PPT (COM) ----------------
    def _convert_ppt(self, path: Path):
        """Convierte .ppt a TXT.

        Preferencia:
        1) PowerPoint vía COM (si hay Office + pywin32).
        2) LibreOffice headless (si está instalado).
        """
        # 1) PowerPoint por COM
        try:
            import pythoncom  # type: ignore
            pythoncom.CoInitialize()
            import win32com.client  # type: ignore

            try:
                ppt = win32com.client.gencache.EnsureDispatch("PowerPoint.Application")
            except Exception:
                ppt = win32com.client.Dispatch("PowerPoint.Application")
            ppt.Visible = False
            pres = ppt.Presentations.Open(str(path), WithWindow=False)

            chunks: List[str] = []
            total = max(1, pres.Slides.Count)
            for i in range(1, pres.Slides.Count + 1):
                slide = pres.Slides(i)
                texts: List[str] = []
                try:
                    for j in range(1, slide.Shapes.Count + 1):
                        shp = slide.Shapes(j)
                        try:
                            if getattr(shp, "HasTextFrame", 0):
                                tf = shp.TextFrame
                                if getattr(tf, "HasText", 0):
                                    texts.append(tf.TextRange.Text or "")
                            if hasattr(shp, "HasTable") and getattr(shp, "HasTable", 0):
                                tbl = shp.Table
                                for r in range(1, tbl.Rows.Count + 1):
                                    row_cells = []
                                    for c in range(1, tbl.Columns.Count + 1):
                                        try:
                                            row_cells.append(tbl.Cell(r, c).Shape.TextFrame.TextRange.Text or "")
                                        except Exception:
                                            row_cells.append("")
                                    if row_cells:
                                        texts.append("\t".join(row_cells))
                        except Exception:
                            pass
                except Exception:
                    pass

                chunks.append(f"[Diapositiva {i}]\n" + "\n".join(texts).strip())
                self._emit("extract_page", i, total)

            out_file = path.with_suffix(".txt")
            out_file.write_text(_norm_txt("\n\n".join(chunks)), encoding="utf-8")
            try:
                pres.Close()
            except Exception:
                pass
            try:
                ppt.Quit()
            except Exception:
                pass
            self._emit("saved", out_file.name)

            try:
                pythoncom.CoUninitialize()
            except Exception:
                pass
            return

        except Exception as e:
            # Liberar COM si llegó a inicializar.
            try:
                import pythoncom  # type: ignore
                try:
                    pythoncom.CoUninitialize()
                except Exception:
                    pass
            except Exception:
                pass
            self._emit("error", f"PowerPoint/COM no disponible, probando LibreOffice… ({e})")

        # 2) Fallback LibreOffice headless
        if self._fallback_ppt_via_libreoffice(path):
            return

        self._emit("error", "Para .ppt necesitas PowerPoint+pywin32 o LibreOffice instalado para conversión headless.")


    # ---------------- EPUB ----------------
    def _convert_epub(self, path: Path):
        if (not _HAS_EPUB) or (epub is None) or (not _HAS_BS4) or (BeautifulSoup is None):
            self._emit("error", "EPUB no disponible: falta ebooklib/lxml o beautifulsoup4 en este intérprete.")
            return
        book = epub.read_epub(str(path))
        docs = [it for it in book.get_items() if isinstance(it, EpubHtml)]
        pages = len(docs)
        text_pages: List[str] = []

        for i, item in enumerate(docs, 1):
            self._emit("extract_page", i, pages)
            try:
                html = item.get_content().decode("utf-8", errors="ignore")
                txt = BeautifulSoup(html, "html.parser").get_text(separator=" ", strip=True)
            except Exception:
                txt = ""
            text_pages.append(txt)

        out_file = path.with_suffix(".txt")
        out_file.write_text("\n\n".join(text_pages), encoding="utf-8")
        self._emit("saved", out_file.name)

    # ---------------- PDF ----------------
    # ---------------- PDF ----------------
    def _convert_pdf(self, path: Path):
        """PDF -> TXT con RAM mínima y OCR inteligente.

        Reglas:
          - Escribe el .txt inmediatamente por página (sin acumular texto en RAM).
          - NO aplica OCR si >=80% de páginas tienen texto (>=min_chars).
          - Si aplica OCR: solo en páginas con texto < min_chars, usando OCRmyPDF con skip_text.
          - OCR se parsea en streaming desde el sidecar y se guarda en cache en disco (sin dict gigante).
          - El progreso lógico siempre avanza por páginas (OCR no suma unidades).
        """
        import os
        import io
        import gc
        import math
        import tempfile
        import contextlib
        import re

        # ---- OCR config
        ocr_cfg = (CONFIG.get("ocr", {}) or {})
        enable_ocr = bool(ocr_cfg.get("enable", True))
        ocr_force = bool(ocr_cfg.get("force", False))
        ocr_min_chars = int(ocr_cfg.get("min_chars", 20))
        ocr_trigger_ratio = float(ocr_cfg.get("trigger_ratio", 0.10))
        ocr_timeout = float(ocr_cfg.get("timeout_sec", ocr_cfg.get("timeout", 60.0)))
        ocr_skip_text = bool(ocr_cfg.get("skip_text", True))
        ocr_dpi = int(ocr_cfg.get("dpi", 300))
        ocr_lang = str(ocr_cfg.get("lang", "spa+eng") or "spa+eng")
        ocr_jobs_cfg = int(ocr_cfg.get("jobs", 0) or 0)
        ocr_whitelist = str(ocr_cfg.get("tessedit_char_whitelist", ocr_cfg.get("whitelist", "") ) or "")
        ocr_whitelist = re.sub(r"[^A-Za-z0-9]", "", ocr_whitelist)
        ocr_sanitize = bool(ocr_cfg.get("sanitize_ocr_output", True))

        # Resolver disponibilidad real de OCR
        has_ocr = False
        ocrmypdf_mod = None
        if enable_ocr or ocr_force:
            try:
                import ocrmypdf as _ocrmypdf  # type: ignore
                ocrmypdf_mod = _ocrmypdf
                has_ocr = True
            except Exception:
                has_ocr = False
        enable_ocr = bool(enable_ocr and has_ocr)
        if ocr_force and not has_ocr:
            # si se fuerza pero no hay ocrmypdf, caer a extracción base
            ocr_force = False

        out_file = path.with_suffix(".txt")

        # ---- Helpers cache base
        def _cache_write(fh, s: str):
            b = (s or "").encode("utf-8", errors="ignore")
            fh.write(len(b).to_bytes(4, "little", signed=False))
            fh.write(b)

        def _cache_iter(cache_path: Path):
            with open(cache_path, "rb") as fh:
                while True:
                    lb = fh.read(4)
                    if not lb:
                        break
                    ln = int.from_bytes(lb, "little", signed=False)
                    if ln <= 0:
                        yield ""
                        continue
                    yield fh.read(ln).decode("utf-8", errors="ignore")

        # ---- Helpers OCR records (p0 + len + bytes)
        def _ocrrec_write(fh, p0: int, s: str):
            b = (s or "").encode("utf-8", errors="ignore")
            fh.write(int(p0).to_bytes(4, "little", signed=False))
            fh.write(len(b).to_bytes(4, "little", signed=False))
            fh.write(b)

        def _ocrrec_iter(cache_path: Path):
            with open(cache_path, "rb") as fh:
                while True:
                    pb = fh.read(4)
                    if not pb:
                        break
                    p0 = int.from_bytes(pb, "little", signed=False)
                    lb = fh.read(4)
                    if not lb:
                        break
                    ln = int.from_bytes(lb, "little", signed=False)
                    if ln <= 0:
                        yield p0, ""
                        continue
                    yield p0, fh.read(ln).decode("utf-8", errors="ignore")

        def _filter_to_whitelist(s: str) -> str:
            if not ocr_sanitize:
                return s or ""
            if not s:
                return ""
            if not ocr_whitelist:
                return s
            allowed = set(ocr_whitelist)
            allowed.update({" ", "\n", "\r", "\t", "\f"})
            return "".join(ch for ch in s if ch in allowed)

        def _build_pages_arg(p0_list):
            if not p0_list:
                return ""
            nums = [int(p0) + 1 for p0 in sorted(set(p0_list))]
            ranges = []
            a = b = nums[0]
            for n in nums[1:]:
                if n == b + 1:
                    b = n
                else:
                    ranges.append(f"{a}-{b}" if a != b else f"{a}")
                    a = b = n
            ranges.append(f"{a}-{b}" if a != b else f"{a}")
            return ",".join(ranges)

        cpu = os.cpu_count() or 4
        def _pick_ocr_jobs() -> int:
            if ocr_jobs_cfg and ocr_jobs_cfg > 0:
                return max(1, min(int(ocr_jobs_cfg), cpu))
            # Cap conservador para evitar picos de RAM
            return max(1, min(4, max(1, cpu - 1)))

        # ---- Streaming de sidecar por \f
        def _iter_sidecar_segments(sidecar_path: Path):
            import os
            CHUNK = 65536
            buf = bytearray()
            with open(sidecar_path, "rb") as fh:
                while True:
                    chunk = fh.read(CHUNK)
                    if not chunk:
                        break
                    buf.extend(chunk)
                    while True:
                        i = buf.find(b"\x0c")
                        if i < 0:
                            break
                        seg = bytes(buf[:i])
                        del buf[:i+1]
                        yield seg.decode("utf-8", errors="ignore")
            if buf:
                yield bytes(buf).decode("utf-8", errors="ignore")

        # ---- Extract base (preferir fitz)
        use_fitz = bool(globals().get('_HAS_FITZ', False) and globals().get('fitz', None) is not None)
        doc = None
        reader = None
        pages = 0
        try:
            if use_fitz:
                doc = fitz.open(str(path))
                pages = int(getattr(doc, 'page_count', 0) or 0)
            else:
                reader = PdfReader(str(path))
                try:
                    pages = int(reader.get_num_pages())
                except Exception:
                    pages = len(reader.pages)
        except Exception as e:
            self._emit('error', f'Falló lectura PDF: {e}')
            try:
                if doc is not None:
                    doc.close()
            except Exception:
                pass
            return

        if pages <= 0:
            try:
                if doc is not None:
                    doc.close()
            except Exception:
                pass
            self._emit('error', 'PDF vacío.')
            return

        # ---- Extracción + escritura inmediata
        short_pages = []  # p0 con texto < min_chars
        textful_count = 0
        cache_path = None
        try:
            with tempfile.NamedTemporaryFile('wb', suffix='.miaucache', delete=False) as tf, \
                 open(out_file, 'w', encoding='utf-8', errors='ignore') as out:
                cache_path = Path(tf.name)
                for p0 in range(pages):
                    try:
                        if use_fitz and doc is not None:
                            txt = doc.load_page(p0).get_text('text') or ''
                        else:
                            txt = self._safe_extract_pdf_page(reader.pages[p0]) if reader is not None else ''
                    except Exception:
                        txt = ''
                    txt = _norm_txt(txt)
                    _cache_write(tf, txt)

                    # escribir inmediatamente
                    out.write(txt or '')
                    if p0 + 1 < pages:
                        out.write('\n')
                    out.flush()

                    st = (txt or '').strip()
                    if len(st) >= ocr_min_chars:
                        textful_count += 1
                    if len(st) < ocr_min_chars:
                        short_pages.append(p0)

                    self._emit('extract_page', p0 + 1, pages)
        finally:
            try:
                if doc is not None:
                    doc.close()
            except Exception:
                pass
            try:
                if use_fitz and globals().get('fitz', None) is not None:
                    try:
                        fitz.TOOLS.store_shrink(100)
                    except Exception:
                        pass
            except Exception:
                pass

        # ---- Decisión OCR con regla >=80% texto
        textful_ratio = (textful_count / float(pages)) if pages else 0.0
        short_ratio = (len(short_pages) / float(pages)) if pages else 0.0
        do_ocr = False
        if ocr_force and enable_ocr:
            do_ocr = True
        elif enable_ocr:
            if textful_ratio >= 0.80:
                do_ocr = False
            else:
                do_ocr = bool(short_ratio > float(ocr_trigger_ratio))

        if (not do_ocr) or (not short_pages) or (ocrmypdf_mod is None):
            # sin OCR, limpiar cache y salir
            try:
                if cache_path is not None:
                    cache_path.unlink(missing_ok=True)
            except Exception:
                pass
            try:
                del reader
            except Exception:
                pass
            gc.collect()
            self._emit('saved', out_file.name)
            return

        # ---- OCR (chunked) => records en disco
        self._emit('ocr_start', len(short_pages))
        ocr_cache_path = None
        try:
            with tempfile.NamedTemporaryFile('wb', suffix='.miaucacheocr', delete=False) as ocf:
                ocr_cache_path = Path(ocf.name)
                chunk_size = int((CONFIG.get('performance', {}) or {}).get('ocr_chunk_size', 50) or 50)
                chunk_size = max(5, min(chunk_size, 200))

                done = 0
                pages_sorted = sorted(set(short_pages))
                for i0 in range(0, len(pages_sorted), chunk_size):
                    chunk = pages_sorted[i0:i0+chunk_size]
                    pages_arg = _build_pages_arg(chunk)
                    if not pages_arg:
                        done = min(len(pages_sorted), i0 + len(chunk))
                        self._emit('ocr_page', done, len(pages_sorted))
                        continue

                    sidecar = None
                    out_pdf = None
                    try:
                        with tempfile.NamedTemporaryFile('wb', suffix='.pdf', delete=False) as tfp:
                            out_pdf = Path(tfp.name)
                        with tempfile.NamedTemporaryFile('wb', suffix='.txt', delete=False) as tfs:
                            sidecar = Path(tfs.name)

                        tcfg = {}
                        if ocr_whitelist:
                            tcfg['tessedit_char_whitelist'] = ocr_whitelist

                        with contextlib.redirect_stdout(io.StringIO()), contextlib.redirect_stderr(io.StringIO()):
                            kwargs = dict(
                                pages=pages_arg,
                                skip_text=ocr_skip_text,
                                force_ocr=False,
                                jobs=int(_pick_ocr_jobs()),
                                tesseract_timeout=float(ocr_timeout),
                                oversample=int(ocr_dpi),
                                language=ocr_lang,
                                output_type='none',
                                sidecar=str(sidecar),
                                tesseract_config=tcfg if tcfg else None,
                            )
                            try:
                                ocrmypdf_mod.ocr(str(path), str(out_pdf), **kwargs)
                            except TypeError:
                                kwargs.pop('output_type', None)
                                try:
                                    ocrmypdf_mod.ocr(str(path), str(out_pdf), **kwargs)
                                except TypeError:
                                    kwargs.pop('tesseract_config', None)
                                    ocrmypdf_mod.ocr(str(path), str(out_pdf), **kwargs)

                        # Parsear sidecar en streaming y escribir records
                        seg_iter = _iter_sidecar_segments(sidecar)
                        for j, p0 in enumerate(chunk):
                            try:
                                seg = next(seg_iter)
                            except StopIteration:
                                seg = ''
                            seg = _filter_to_whitelist(_norm_txt(seg)).strip('\n')
                            _ocrrec_write(ocf, int(p0), seg)

                    except Exception:
                        # si falla OCR chunk, escribir vacío para esas páginas
                        for p0 in chunk:
                            _ocrrec_write(ocf, int(p0), '')
                    finally:
                        for tmp in (sidecar, out_pdf):
                            if tmp is not None:
                                try:
                                    tmp.unlink()
                                except Exception:
                                    pass

                    done = min(len(pages_sorted), i0 + len(chunk))
                    self._emit('ocr_page', done, len(pages_sorted))

            # ---- Re-ensamble final (sin RAM)
            tmp_out = None
            try:
                with tempfile.NamedTemporaryFile('w', suffix='.txt', delete=False, encoding='utf-8', errors='ignore') as tfout:
                    tmp_out = Path(tfout.name)

                ocr_it = _ocrrec_iter(ocr_cache_path)
                try:
                    cur = next(ocr_it)
                except StopIteration:
                    cur = None

                with open(tmp_out, 'w', encoding='utf-8', errors='ignore') as out:
                    for p0, base_txt in enumerate(_cache_iter(cache_path)):
                        base_txt = base_txt or ''
                        ocr_txt = ''
                        if cur is not None and int(cur[0]) == int(p0):
                            ocr_txt = cur[1] or ''
                            try:
                                cur = next(ocr_it)
                            except StopIteration:
                                cur = None
                        use_txt = ocr_txt if len((ocr_txt or '').strip()) > len((base_txt or '').strip()) else base_txt
                        out.write(use_txt or '')
                        if p0 + 1 < pages:
                            out.write('\n')

                os.replace(str(tmp_out), str(out_file))
            finally:
                if tmp_out is not None:
                    try:
                        tmp_out.unlink(missing_ok=True)
                    except Exception:
                        pass

        finally:
            for tmp in (cache_path, ocr_cache_path):
                if tmp is not None:
                    try:
                        tmp.unlink(missing_ok=True)
                    except Exception:
                        pass
            try:
                del reader
            except Exception:
                pass
            gc.collect()

        self._emit('saved', out_file.name)



    def _write_simple(self, text: str, path: Path):
        out_file = path.with_suffix(".txt")
        out_file.write_text(text or "", encoding="utf-8")
        self._emit("saved", out_file.name)

    def _process_files(self):
        total_files = len(self._files)
        # Total global estimado (para que el % sea de TODA la lista)
        try:
            total_units = sum(_estimate_units_for_path(p) for p in self._files)
        except Exception:
            total_units = max(1, total_files)
        self._emit("global_pages_total", int(max(1, total_units)))
        self._emit("global_total", int(max(1, total_units)))
        for idx, path in enumerate(self._files, 1):
            try:
                self._emit("file", idx, total_files, path.name)
                if not path.is_file():
                    self._emit("error", f"No existe: {path.name}")
                    continue
                sfx = path.suffix.lower()
                if sfx == ".pdf":
                    self._convert_pdf(path)
                elif sfx == ".epub":
                    self._convert_epub(path)
                elif sfx == ".docx":
                    self._convert_docx(path)
                elif sfx == ".pptx":
                    self._convert_pptx(path)
                elif sfx == ".doc":
                    self._convert_doc(path)
                elif sfx == ".ppt":
                    self._convert_ppt(path)
                else:
                    self._emit("error", f"No soportado: {sfx}")
            except Exception as e:
                self._emit("error", f"Falló {path.name}: {e}")
        self._emit("done")

    # ---------------- UI loop ----------------
    def _pump_progress(self):
        """Bombea la cola de progreso sin bloquear la UI.

        Problema previo: drenar la cola completa en un solo callback puede impedir el repaint
        cuando el worker produce eventos continuamente. Se limita por tiempo/cantidad.
        """
        import time as _time

        max_events = 250
        max_ms = 12.0
        start = _time.perf_counter()
        processed = 0

        try:
            while processed < max_events and ((_time.perf_counter() - start) * 1000.0) < max_ms:
                try:
                    tag, *data = self._queue.get_nowait()
                except queue.Empty:
                    break

                processed += 1

                if tag == "global_total":
                    self._g_total_units = int(max(1, data[0] if data else 1))
                    self._pdlg.set_global_span(self._g_total_units)
                    self._pdlg.set_done(self._g_done_units)
                    continue

                if tag == "global_pages_total":
                    self._g_pages_total = int(max(1, data[0] if data else 1))
                    # reset lógico
                    if self._g_pages_done > self._g_pages_total:
                        self._g_pages_done = self._g_pages_total
                    continue

                if tag == "file":
                    idx, total, name = data
                    self._cur_file_name = str(name)
                    self._last_file_unit = 0
                    self._last_ocr_unit = 0
                    self._pdlg.set_title(f"Procesando ({idx}/{total}): {name}")
                    self._pdlg.set_subtitle("Preparando…")
                    continue

                if tag == "extract_page":
                    i, total = data
                    i = int(i)
                    total = int(max(1, total))
                    delta = i - int(self._last_file_unit)
                    if delta < 0:
                        delta = i
                    self._last_file_unit = i

                    if delta > 0:
                        self._g_done_units += delta
                        self._g_pages_done += delta
                        self._pdlg.set_done(self._g_done_units)

                    self._pdlg.set_subtitle(f"Página {i}/{total}")
                    continue

                if tag == "ocr_start":
                    pages_ocr = int(max(1, data[0]))
                    self._last_ocr_unit = 0
                    self._pdlg.set_subtitle(f"OCR 0/{pages_ocr}")
                    continue

                if tag == "ocr_page":
                    i, total = data
                    i = int(i)
                    total = int(max(1, total))
                    delta = i - int(self._last_ocr_unit)
                    if delta < 0:
                        delta = i
                    self._last_ocr_unit = i
                    self._pdlg.set_subtitle(f"OCR {i}/{total}")
                    continue

                if tag == "saved":
                    self._pdlg.set_subtitle(f"Guardado: {data[0]}")
                    continue

                if tag == "error":
                    self._pdlg.set_subtitle(f"Error: {data[0]}")
                    continue

                if tag == "done":
                    self._pdlg.set_done(self._g_total_units or 1)
                    self._pdlg.set_progress(1.0)
                    self._pdlg.set_title("Listo")
                    self._pdlg.set_subtitle("Listo")
                    # Mantener visible ~2s y luego cerrar todo
                    try:
                        delay_ms = int((get_config("progress") or {}).get("auto_close_ms", 2000))
                    except Exception:
                        delay_ms = 2000
                    try:
                        self._root.after(delay_ms, self._final_close)
                    except Exception:
                        try:
                            self._final_close()
                        except Exception:
                            pass
                    return

        except Exception:
            pass

        self._root.after(30, self._pump_progress)


    def _final_close(self):
        """Cierre robusto: asegura terminar mainloop y destruir la ventana."""
        try:
            try:
                self._root.quit()
            except Exception:
                pass
            self._graceful_close()
        except Exception:
            try:
                self._root.destroy()
            except Exception:
                pass

    def _graceful_close(self):
        try:
            if getattr(self, "_pdlg", None) and self._pdlg.winfo_exists():
                try:
                    self._pdlg.close()
                except Exception:
                    self._pdlg.destroy()
        finally:
            try:
                self._root.destroy()
            except Exception:
                pass
            try:
                self._root.quit()
            except Exception:
                pass


# =============================================================================
# ──────────────────────────────────────────────────────────────────────────────
def _merge_multiselect_startup(
    files_in,
    *,
    mutex_name: str,
    spool_name: str,
    settle_ms: int = 7000,
    stable_ms: int = 650,
    poll_ms: int = 60,
):
    """
    Agrupa multi-selección de Explorer cuando éste ejecuta 1 proceso por archivo.
    - Si ya existe otra instancia (mutex), este proceso solo APPENDEA sus rutas a un spool y SALE.
    - La instancia "primaria" colecta rutas desde el spool hasta que el set se estabiliza.
    Mantiene el arranque ágil: si solo hay 1 archivo, normalmente se estabiliza en <200 ms.
    """
    import os
    from pathlib import Path
    from typing import List

    if os.name != "nt":
        return list(files_in), True

    try:
        import time
        import tempfile
        import ctypes
        import ctypes.wintypes as _wt
        import msvcrt

        kernel32 = ctypes.WinDLL("kernel32", use_last_error=True)
        CreateMutexW = kernel32.CreateMutexW
        CreateMutexW.argtypes = [_wt.LPVOID, _wt.BOOL, _wt.LPCWSTR]
        CreateMutexW.restype = _wt.HANDLE
        GetLastError = kernel32.GetLastError
        CloseHandle = kernel32.CloseHandle

        ERROR_ALREADY_EXISTS = 183
        mutex = CreateMutexW(None, False, mutex_name)
        already = (GetLastError() == ERROR_ALREADY_EXISTS)

        base = os.environ.get("LOCALAPPDATA") or os.environ.get("APPDATA") or tempfile.gettempdir()
        spool_dir = Path(base) / "MiausoftSuite" / "IPC"
        spool_dir.mkdir(parents=True, exist_ok=True)
        spool = spool_dir / spool_name

        def _append(paths: List[str]) -> None:
            if not paths:
                return
            with spool.open("a+", encoding="utf-8") as f:
                try:
                    # lock 1 byte region (suficiente para exclusión mutua entre procesos)
                    msvcrt.locking(f.fileno(), msvcrt.LK_LOCK, 1)
                except Exception:
                    pass
                f.seek(0, os.SEEK_END)
                for p in paths:
                    p = (p or "").replace("\n", "").strip()
                    if p:
                        f.write(p + "\n")
                f.flush()
                try:
                    msvcrt.locking(f.fileno(), msvcrt.LK_UNLCK, 1)
                except Exception:
                    pass

        def _drain() -> List[str]:
            out: List[str] = []
            with spool.open("a+", encoding="utf-8") as f:
                try:
                    msvcrt.locking(f.fileno(), msvcrt.LK_LOCK, 1)
                except Exception:
                    pass
                f.seek(0)
                out = [(line or "").strip() for line in f]
                try:
                    f.seek(0)
                    f.truncate()
                except Exception:
                    pass
                try:
                    msvcrt.locking(f.fileno(), msvcrt.LK_UNLCK, 1)
                except Exception:
                    pass
            return [p for p in out if p]

        files_norm = []
        for x in list(files_in):
            try:
                files_norm.append(str(x))
            except Exception:
                pass

        if not files_norm:
            return list(files_in), True

        # Preferir tiempos desde config central (sin configs locales)
        try:
            from config import get_config as _get_cfg  # centralizado
            _cfg = _get_cfg(None) if callable(_get_cfg) else {}
            _ms = ((_cfg or {}).get("app", {}) or {}).get("multiselect", {}) or {}
            settle_ms = int(_ms.get("settle_ms", settle_ms))
            stable_ms = int(_ms.get("stable_ms", stable_ms))
            poll_ms = int(_ms.get("poll_ms", poll_ms))
        except Exception:
            pass

        # Si ya llegó más de 1 ruta en un mismo proceso, no hay que esperar casi nada.
        if len(files_norm) > 1:
            stable_ms = min(stable_ms, 220)
            settle_ms = min(settle_ms, 1200)


        if already:
            _append(files_norm)
            try:
                CloseHandle(mutex)
            except Exception:
                pass
            return [], False

        # Primaria: escribe sus rutas y colecta las demás hasta estabilizar.
        _append(files_norm)

        merged: List[str] = []
        seen: set[str] = set()
        stable = 0
        deadline = time.monotonic() + (max(200, int(settle_ms)) / 1000.0)

        while time.monotonic() < deadline and stable < int(stable_ms):
            new = _drain()
            added = 0
            for p in new:
                if not p or p in seen:
                    continue
                try:
                    if Path(p).exists():
                        seen.add(p)
                        merged.append(p)
                        added += 1
                except Exception:
                    continue
            if added == 0:
                stable += int(poll_ms)
            else:
                stable = 0
            time.sleep(max(10, int(poll_ms)) / 1000.0)

        # Drenaje final por si llegó algo entre el último sleep y el cierre.
        for p in _drain():
            if p and p not in seen:
                try:
                    if Path(p).exists():
                        seen.add(p)
                        merged.append(p)
                except Exception:
                    pass

        try:
            spool.unlink(missing_ok=True)
        except Exception:
            pass
        # Intencionalmente NO cerramos el mutex aquí.
        # Motivo: Explorer puede lanzar 1 proceso por archivo seleccionado;
        # si liberamos el mutex demasiado pronto, otra instancia puede volverse
        # "primaria" y abrir otra ventana. El sistema cerrará el handle al
        # terminar el proceso, y el objeto se destruye cuando se cierra el último handle.

        # Si por alguna razón quedó vacío, conservamos el input original.
        return (merged if merged else files_norm), True

    except Exception:
        return list(files_in), True

def _safe_main(argv):
    import os, traceback
    try:
        raw_args = list(argv[1:])
        files: List[str] = []
        for a in raw_args:
            if not a or a.startswith("-"):
                continue
            s = a.strip().strip('"').strip("'")
            s = os.path.expandvars(s)
            try:
                from pathlib import Path as _P
                pp = _P(s).expanduser().resolve()
            except Exception:
                pp = _P(s)
            if pp.exists() and pp.is_file():
                files.append(str(pp))


        # Agrupar multi-selección cuando Explorer ejecuta 1 instancia POR archivo seleccionado.


        files, _launch = _merge_multiselect_startup(


            files,


            mutex_name=r"Local\MiausoftSuite_TCompleto_MultiSelect",


            spool_name="tcompleto.paths.txt",


        )


        if not _launch:


            return



        HeadlessConverterApp(files)
    except Exception:
        try:
            # Log SIEMPRE en un lugar escribible (evita Program Files / permisos).
            base = os.environ.get("LOCALAPPDATA") or os.environ.get("APPDATA") or str(Path.home())
            log_dir = Path(base) / "MiausoftSuite" / "Logs"
            log_dir.mkdir(parents=True, exist_ok=True)
            log_path = log_dir / "tcompleto.error.log"
            log_path.write_text(traceback.format_exc(), encoding="utf-8")
        except Exception:
            log_path = None
        try:
            import tkinter as _tk
            from tkinter import messagebox as _mb
            root = _tk.Tk(); root.withdraw()
            msg = "Error al iniciar el conversor."
            if log_path:
                msg += f"\n\nRevisa el log:\n{log_path}"
            _mb.showerror("MiausoftSuite", msg)
        except Exception:
            print("Error al iniciar:", traceback.format_exc())


if __name__ == "__main__":
    _safe_main(sys.argv)